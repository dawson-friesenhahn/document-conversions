import pptx
import pptx.shapes
import pptx.slide
import pptx.table
import pptx.text.text
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE

from PIL import Image, ImageDraw, ImageFont
import os
import io

from reportlab.pdfgen import canvas
import glob
import shutil

from .util import emu_to_pixels, emu_to_pt, pt_to_emu, create_slide_number_string
from . import ARIAL_FONT


def render_slide_as_image(slide: pptx.slide.Slide, slide_width_emu, slide_height_emu):
    # Convert slide dimensions from EMUs to pixels
    slide_width = emu_to_pixels(slide_width_emu)
    slide_height = emu_to_pixels(slide_height_emu)

    # Create an image canvas (in RGB mode)
    image = Image.new("RGB", (slide_width, slide_height), color=(255, 255, 255))
    draw = ImageDraw.Draw(image)

    # Loop through each shape on the slide and render its content
    for shape in slide.shapes:
        # if shape.is_placeholder:
        #     continue  # Skip placeholders
        if shape.has_text_frame:
            y_offset: int = 0
            # Render text
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    
                    # Convert the left/top position from EMUs to pixels
                    x = emu_to_pixels(shape.left)
                    y = emu_to_pixels(shape.top + y_offset)

                    if run.font.size:
                        y_offset += run.font.size
                        font_size_pt = emu_to_pt(run.font.size)
                    else:
                        y_offset += pt_to_emu(12)
                        font_size_pt= 12

                    # print(f"({x},{y}): {run.text}")
                    #font= ImageFont.load_default(size=font_size_pt)                   
                    font= ImageFont.truetype(ARIAL_FONT, size=font_size_pt)
                    # Draw text onto the image (adjust position and styling as needed)
                    try:
                        draw.text(
                            (x, y),
                            run.text,
                            font=font,
                            fill=(0, 0, 0),
                        )
                    except Exception as e:
                        print(e)

        # Example for simple rectangles (can be extended to other shapes)
        if False and shape.shape_type == MSO_AUTO_SHAPE_TYPE.RECTANGLE:  # Rectangle
            x1 = emu_to_pixels(shape.left)
            y1 = emu_to_pixels(shape.top)
            x2 = x1 + emu_to_pixels(shape.width)
            y2 = y1 + emu_to_pixels(shape.height)

            draw.rectangle([x1, y1, x2, y2], outline="black", width=3)

        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table: pptx.table.Table = shape.table

            # Get table dimensions
            rows = len(table.rows)
            cols = len(table.columns)

            # Calculate the cell width and height
            table_width = emu_to_pixels(shape.width)
            table_height = emu_to_pixels(shape.height)
            row_height = table_height // rows
            col_width = (
                table_width // cols
            )  # TODO this isn't correct, the columns/rows might not have uniform width or height

            # Start drawing the table
            for i, row in enumerate(table.rows):
                for j, cell in enumerate(row.cells):
                    # Calculate the cell position
                    x1 = emu_to_pixels(shape.left) + j * col_width
                    y1 = emu_to_pixels(shape.top) + i * row_height
                    x2 = x1 + col_width
                    y2 = y1 + row_height

                    # Draw the rectangle for the cell
                    draw.rectangle([x1, y1, x2, y2], outline="black", width=2)

                    # Get the cell text and render it
                    cell_text = cell.text.strip()
                    if cell_text:
                        min_font_size = 1e9
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.size < min_font_size:
                                    min_font_size = run.font.size
                        
                        if min_font_size > 1000:
                            min_font_size = 12

                        # Font size for table text (adjust as needed)
                        font_size_pt = emu_to_pt(
                            min_font_size
                        )  # TODO dynamically get the font size

                        if not font_size_pt:
                            font_size_pt = 12
                        #font = ImageFont.load_default(size=font_size_pt)
                        font= ImageFont.truetype(ARIAL_FONT, size=font_size_pt)

                        # Calculate the bounding box of the text
                        bbox = draw.textbbox((0, 0), cell_text, font=font)
                        text_width = bbox[2] - bbox[0]  # width of the text
                        text_height = bbox[3] - bbox[1]  # height of the text

                        # Center the text inside the cell
                        text_x = x1 + (col_width - text_width) // 2
                        text_y = y1 + (row_height - text_height) // 2
                        draw.text(
                            (text_x, text_y), cell_text, font=font, fill=(0, 0, 0)
                        )

        # Handle Picture shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # Picture
            # Extract the picture and convert to an image object
            image_stream = io.BytesIO(shape.image.blob)
            img = Image.open(image_stream)

            # Convert the image position and size from EMUs to pixels
            x1 = emu_to_pixels(shape.left)
            y1 = emu_to_pixels(shape.top)
            width = emu_to_pixels(shape.width)
            height = emu_to_pixels(shape.height)

            # Resize the image if necessary (it may need to be resized to fit the slide)
            img = img.resize((width, height))

            # Paste the image onto the slide image
            image.paste(img, (x1, y1))

    return image


def pptx_as_images(pptx_path) -> list[Image.Image]:
    prs = pptx.Presentation(pptx_path)

    # Get the slide dimensions from the presentation object
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height

    ret = []

    for slide_num, slide in enumerate(prs.slides):
        # Render the slide as an image
        image = render_slide_as_image(slide, slide_width_emu, slide_height_emu)

        ret.append(image)
    return ret


def save_pptx_as_images(pptx_path, output_folder) -> str:
    prs = pptx.Presentation(pptx_path)
    os.makedirs(output_folder, exist_ok=True)
    output_folder = os.path.join(output_folder, "temp_images")

    if os.path.exists(output_folder):
        shutil.rmtree(output_folder)

    os.mkdir(output_folder)

    # Get the slide dimensions from the presentation object
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height

    for slide_num, slide in enumerate(prs.slides):
        # Render the slide as an image
        image = render_slide_as_image(slide, slide_width_emu, slide_height_emu)

        # Save the image as PNG
        image.save(f"{output_folder}/slide_{create_slide_number_string(slide_num+1)}.png")

    return output_folder


def images_to_pdf(images_dir: str, filename: str):
    c = canvas.Canvas(filename=filename, pagesize=(841.89, 595.27))
    images = glob.glob(os.path.join(images_dir, "*.png"))
    for image in images:
        c.drawImage(
            image, 0, 0, width=841.89, height=595.27
        )  # preserveAspectRatio=True)
        c.showPage()

    c.save()


def pptx_to_pdf(pptx_path, output_folder=None, delete_intermediate_images=True) -> str:
    output_name = os.path.basename(pptx_path).split(".")[0] + ".pdf"
    if not output_folder:
        output_folder = os.getcwd()
    if output_folder:
        output_name = os.path.join(output_folder, output_name)

    images_dir = save_pptx_as_images(
        pptx_path, output_folder
    )  # will create new directory there and save the images

    # print(slides_as_images)

    images_to_pdf(images_dir, output_name)

    if delete_intermediate_images:
        shutil.rmtree(images_dir)

    return output_name
